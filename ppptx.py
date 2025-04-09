from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

# Create a presentation object
prs = Presentation()

# Set slide dimensions to 16:9 aspect ratio
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

# Function to add a background image to a slide
def add_background_image(slide, image_path):
    left = top = Inches(0)
    width = prs.slide_width
    height = prs.slide_height
    slide.shapes.add_picture(image_path, left, top, width, height)

# Function to add a semi-transparent overlay
def add_overlay(slide, rgb_color, transparency):
    left = top = Inches(0)
    width = prs.slide_width
    height = prs.slide_height
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*rgb_color)
    shape.fill.transparency = transparency

# Function to add slides with a background image and overlay
def add_slide(title_text, content_text, bg_image_path, overlay_color=(0, 0, 0), overlay_transparency=0.3):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    # Add background image
    add_background_image(slide, bg_image_path)
    
    # Add semi-transparent overlay
    add_overlay(slide, overlay_color, overlay_transparency)
    
    # Add title and content
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = title_text
    content.text = content_text
    
    # Set Title Font Size, Color, and Alignment
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # White
    title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Set Content Font Size, Color, and Alignment
    for paragraph in content.text_frame.paragraphs:
        paragraph.font.size = Pt(24)
        paragraph.font.color.rgb = RGBColor(255, 255, 255)  # White
        paragraph.alignment = PP_ALIGN.LEFT

# Title Slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
add_background_image(slide, 'D:\\code\\bgm.jpg')
add_overlay(slide, (0, 0, 0), 0.5)  # Dark overlay for better text visibility

title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Credit Card Fraud Detection Using Machine Learning"
subtitle.text = "Chinmay Malve\nIIIT Allahabad, June-July 2024"

title.text_frame.paragraphs[0].font.size = Pt(54)
title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
subtitle.text_frame.paragraphs[0].font.size = Pt(32)
subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

# Adding content slides with background images
add_slide("Introduction", 
          "• Credit card fraud has become a significant concern with the rise of digital payments.\n"
          "• Detecting fraudulent transactions is crucial to protecting consumers and financial institutions.\n"
          "• Objective: Build a model to detect 100% of fraudulent transactions while minimizing false positives.", 
          'D:\\code\\bgm.jpg')

add_slide("Problem Statement", 
          "• Highly imbalanced dataset with only 0.172% fraudulent transactions.\n"
          "• Challenge: Detect fraud accurately without falsely flagging legitimate transactions.\n"
          "• Fraud detection involves monitoring user behavior and using ML models to classify transactions.", 
          'D:\\code\\bgm.jpg')

add_slide("Dataset Overview", 
          "• Source: Kaggle (credit card transactions from European cardholders, 2013).\n"
          "• 284,807 transactions, 492 fraud cases (0.172%).\n"
          "• Features: Time, Amount, 28 PCA-transformed features (V1-V28), Class.", 
          'D:\\code\\bgm.jpg')

add_slide("Methodology", 
          "1. Data Collection\n"
          "2. Data Preprocessing (Standardized Time & Amount)\n"
          "3. Model Implementation (LOF, Isolation Forest)\n"
          "4. Model Testing & Evaluation\n"
          "Tools: Python (Jupyter Notebook), sklearn, NumPy, SciPy, Matplotlib", 
          'D:\\code\\bgm.jpg')

add_slide("Data Preparation", 
          "• Standardized 'Time' and 'Amount'.\n"
          "• Imbalanced data was handled as is, without oversampling.\n"
          "• Exploratory Data Analysis (EDA) revealed distribution differences between fraud and non-fraud transactions.", 
          'D:\\code\\bgm.jpg')

add_slide("Exploratory Data Analysis (EDA)", 
          "• Visualized transaction volumes across time: Fewer transactions at night.\n"
          "• Fraud transactions have a wider range of amounts.\n"
          "• Correlation heatmaps show strong relations between features.", 
          'D:\\code\\bgm.jpg')

add_slide("Anomaly Detection Models", 
          "• Local Outlier Factor (LOF): Detects anomalies based on local density deviations.\n"
          "• Isolation Forest: Randomly partitions data to isolate anomalies.", 
          'D:\\code\\bgm.jpg')

add_slide("Local Outlier Factor (LOF)", 
          "• LOF calculates anomaly scores by comparing local density with neighbors.\n"
          "• Detected anomalies with moderate accuracy.", 
          'D:\\code\\bgm.jpg')

add_slide("Isolation Forest", 
          "• Partitions features randomly to isolate potential anomalies.\n"
          "• Achieved 97% accuracy with the full dataset.", 
          'D:\\code\\bgm.jpg')

add_slide("Comparison of Models", 
          "• Local Outlier Factor: Moderate accuracy.\n"
          "• Isolation Forest: Higher precision, handled imbalanced data better.\n"
          "• Conclusion: Isolation Forest performed better in both accuracy and efficiency.", 
          'D:\\code\\bgm.jpg')

add_slide("Evaluation Metrics", 
          "• Precision-Recall Curve used due to class imbalance.\n"
          "• Confusion matrix focused on false positives and true positives.", 
          'D:\\code\\bgm.jpg')

add_slide("Test Results", 
          "• LOF: Moderate precision, lower recall.\n"
          "• Isolation Forest: Achieved 97% accuracy on the full dataset.", 
          'D:\\code\\bgm.jpg')

add_slide("Conclusion", 
          "• Fraud detection is crucial in financial systems.\n"
          "• Isolation Forest provided the best results for detecting fraudulent transactions.\n"
          "• Machine learning greatly improves the speed and accuracy of fraud detection.", 
          'D:\\code\\bgm.jpg')

add_slide("Future Scope", 
          "• Tune Random Forest and other ensemble methods.\n"
          "• Use non-anonymized datasets to identify specific factors in fraud detection.\n"
          "• Implement real-time fraud detection for practical use.", 
          'D:\\code\\bgm.jpg')

add_slide("References", 
          "1. Credit Card Fraud Detection Based on Transaction Behavior - IEEE TENCON 2017.\n"
          "2. Visualizing High-Dimensional Data Using t-SNE - Journal of Machine Learning Research (2014).\n"
          "3. Kaggle Credit Card Fraud Detection Dataset (2018).", 
          'D:\\code\\bgm.jpg')

# Save the presentation
prs.save("Credit_Card_Fraud_Detection_Presentation_Enhanced.pptx")

print("Enhanced presentation created successfully!")